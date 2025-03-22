import json
import streamlit as st
from typing import TypedDict
from langgraph.graph import StateGraph, END
import openai
import os
from dotenv import load_dotenv
from io import BytesIO
from docx import Document
from pptx import Presentation
import pdfplumber
import pandas as pd
from bs4 import BeautifulSoup
import mimetypes
import tempfile
import pypandoc

# Load environment variables
load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")

# Define a state schema
class DocumentState(TypedDict):
    uploaded_file: BytesIO
    file_type: str
    summary: str
    chunks: list
    message: str

# Initialize LangGraph with a state schema
upload_graph = StateGraph(state_schema=DocumentState)

# File extraction functions

def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(BytesIO(file.read())) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text.strip()

def extract_text_from_doc(file, file_type):
    if file_type == "docx":
        doc = Document(BytesIO(file.read()))
        return "\n".join([para.text for para in doc.paragraphs])
    elif file_type == "doc":
        with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as temp_doc:
            temp_doc.write(file.read())
            temp_doc_path = temp_doc.name
        temp_docx_path = temp_doc_path + "x"
        try:
            pypandoc.convert_file(temp_doc_path, "docx", outputfile=temp_docx_path)
            doc = Document(temp_docx_path)
            return "\n".join([para.text for para in doc.paragraphs])
        finally:
            os.remove(temp_doc_path)
            if os.path.exists(temp_docx_path):
                os.remove(temp_docx_path)
    else:
        return ""

def extract_text_from_pptx(file):
    prs = Presentation(BytesIO(file.read()))
    return "\n".join([
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    ]).strip()

def extract_text_from_txt(file):
    return file.read().decode("utf-8").strip()

def extract_text_from_html(file):
    soup = BeautifulSoup(file, "html.parser")
    return soup.get_text().strip()

def extract_text_from_any_file(uploaded_file):
    file_type, _ = mimetypes.guess_type(uploaded_file.name)
    
    extractors = {
        "application/pdf": extract_text_from_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": extract_text_from_doc,
        "application/msword": extract_text_from_doc,
        "text/plain": extract_text_from_txt,
        "text/html": extract_text_from_html,
        "application/vnd.ms-powerpoint": extract_text_from_pptx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": extract_text_from_pptx,
        "application/json": lambda file: json.load(file),
        "text/csv": lambda file: pd.read_csv(file).to_string(),
    }
    
    if file_type in extractors:
        return extractors[file_type](uploaded_file, uploaded_file.name.split('.')[-1].lower())
    return "Unsupported file type. Please upload a document or readable file."

def chunk_text(text, chunk_size=500):
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

def summarize_large_text(text, chunk_size=2000):
    chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
    summaries = []
    for chunk in chunks:
        response = openai.ChatCompletion.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"Summarize the following document:\n\n{chunk}"}
            ],
            max_tokens=500
        )
        summaries.append(response['choices'][0]['message']['content'].strip())
    return " ".join(summaries)

# Define LangGraph nodes
def upload_document(state: DocumentState):
    uploaded_file = state["uploaded_file"]
    text = extract_text_from_any_file(uploaded_file)
    if not text:
        return {"summary": "", "chunks": [], "message": "No text extracted from the document."}
    chunks = chunk_text(text)
    summary = summarize_large_text(text)
    return {"summary": summary, "chunks": chunks, "message": "Document processed successfully!"}

# Add nodes to the upload graph
upload_graph.add_node("upload_document", upload_document)
upload_graph.set_entry_point("upload_document")
upload_graph.add_edge("upload_document", END)
upload_app = upload_graph.compile()

def main():
    st.title("Document Upload and Summarization System with LangGraph")
    
    # Document upload and processing
    uploaded_file = st.file_uploader("Upload a Document", type=None)
    if uploaded_file:
        state = {"uploaded_file": uploaded_file}
        result = upload_app.invoke(state)
        st.subheader("Document Summary")
        st.write(result["summary"])

if __name__ == "__main__":
    main()
